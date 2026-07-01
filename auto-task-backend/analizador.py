# analizador.py - CORREGIDO (Gemini con chat nativo, Groq/Cohere truncados, Flujo secuencial reparado)
import os
import re
import time
import shutil
import subprocess
import platform
# import pythoncom

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

from typing import Optional, List, Dict, Any, Callable

from pypdf import PdfReader
from google import genai
from google.genai import types
import groq
import cohere

from motor_docx import GeneradorWord
from generadores.estadistica import GeneradorEstadistica
from config_tareas import CONFIG_TAREAS, TIPO_DEFECTO

# PowerPoint requirements
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


# ══════════════════════════════════════════════════════════════════════════════
# 1. UTILIDADES BÁSICAS
# ══════════════════════════════════════════════════════════════════════════════

MIME_TYPES = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".gif": "image/gif",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
}

def leer_bytes(ruta: str) -> tuple[bytes, str]:
    ext = os.path.splitext(ruta)[1].lower()
    mime = MIME_TYPES.get(ext, "application/octet-stream")
    with open(ruta, "rb") as f:
        return f.read(), mime

def extraer_texto_pdf(ruta: str) -> str:
    try:
        lector = PdfReader(ruta)
        return "\n".join(p.extract_text() for p in lector.pages if p.extract_text())
    except Exception as e:
        return f"[ERROR_PDF: {e}]"

def extraer_texto_de_pdfs(rutas: list) -> str:
    texto_completo = ""
    for r in rutas:
        if r.lower().endswith('.pdf'):
            txt = extraer_texto_pdf(r)
            if not txt.startswith("[ERROR]"):
                texto_completo += f"\n--- Archivo: {os.path.basename(r)} ---\n{txt}\n"
    return texto_completo

def _log(msg: str, cb: Optional[Callable] = None):
    print(msg)
    if cb:
        cb(msg)
### CONVERTIR WORD A PDF (solo Windows con Word instalado) 
def convertir_docx_a_pdf(ruta_docx: str, carpeta_destino: str) -> str | None:
    """
    Convierte un archivo .docx/.doc a .pdf de forma multiplataforma.
    Usa LibreOffice de forma nativa en Linux/Docker, y win32com de respaldo en Windows.
    """
    import os
    import subprocess
    import platform
    import shutil

    if not os.path.exists(ruta_docx):
        _log(f"❌ El archivo {ruta_docx} no existe.", None)
        return None

    ruta_abs = os.path.abspath(ruta_docx)
    carpeta_abs = os.path.abspath(carpeta_destino)
    nombre_base = os.path.splitext(os.path.basename(ruta_docx))[0]
    ruta_pdf = os.path.join(carpeta_abs, nombre_base + ".pdf")

    # 1. EN LINUX/DOCKER (o Windows si LibreOffice está en PATH): Usamos LibreOffice headless
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        try:
            _log("📄 Convirtiendo Word a PDF usando LibreOffice headless...", None)
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", ruta_abs, "--outdir", carpeta_abs],
                capture_output=True,
                timeout=60,
                check=True
            )
            if os.path.exists(ruta_pdf):
                _log(f"✅ Word convertido a PDF (LibreOffice): {ruta_pdf}", None)
                return ruta_pdf
        except Exception as e:
            _log(f"⚠️ Conversión con LibreOffice falló: {e}", None)

    # 2. EN WINDOWS (si no se usó o falló LibreOffice): Usamos win32com con importación local segura
    if platform.system() == "Windows":
        try:
            import pythoncom
            import win32com.client
            
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0 

            doc = None
            try:
                doc = word.Documents.Open(
                    ruta_abs, 
                    ConfirmConversions=False, 
                    ReadOnly=False, 
                    AddToRecentFiles=False
                )
            except Exception:
                try:
                    pv_win = word.ProtectedViewWindows.Open(ruta_abs)
                    doc = pv_win.Edit()
                except Exception as e_pv:
                    raise e_pv

            if doc:
                doc.SaveAs(ruta_pdf, FileFormat=17) # 17 = wdFormatPDF
                doc.Close(SaveChanges=False)
                if os.path.exists(ruta_pdf):
                    _log(f"✅ Word convertido a PDF (Word Windows): {ruta_pdf}", None)
                    return ruta_pdf
        except Exception as e:
            _log(f"❌ Error win32com al convertir en Windows: {e}", None)
        finally:
            if 'word' in locals() and word:
                try:
                    for pv_win in list(word.ProtectedViewWindows):
                        pv_win.Close()
                    word.Quit()
                except:
                    pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass

    return None
### CONVERTIR PDF A WORD (solo Windows con Word instalado)
### CONVERTIR IMAGEN A PDF
def convertir_imagen_a_pdf(ruta_imagen: str, carpeta_destino: str) -> str | None:
    """
    Convierte de forma segura una imagen (PNG, JPG, JPEG, WEBP, GIF) a PDF usando Pillow.
    Maneja canales alfa (transparencia) sustituyéndolos por un fondo blanco para legibilidad.
    """
    from PIL import Image
    import os

    if not os.path.exists(ruta_imagen):
        return None

    try:
        nombre_base = os.path.splitext(os.path.basename(ruta_imagen))[0]
        ruta_pdf = os.path.join(os.path.abspath(carpeta_destino), nombre_base + "_img.pdf")

        with Image.open(ruta_imagen) as img:
            # Si la imagen tiene transparencia (RGBA, LA) o paleta con transparencia (P)
            if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info):
                # Creamos un fondo blanco del mismo tamaño
                fondo_blanco = Image.new("RGB", img.size, (255, 255, 255))
                # Si es RGBA, usamos el canal Alpha como máscara de recorte
                if img.mode == "RGBA":
                    fondo_blanco.paste(img, mask=img.split()[3])
                else:
                    fondo_blanco.paste(img)
                img_final = fondo_blanco
            else:
                # Si no tiene transparencia, conversión directa a RGB
                img_final = img.convert("RGB")
            
            # Guardar como PDF de una sola página
            img_final.save(ruta_pdf, "PDF")

        if os.path.exists(ruta_pdf):
            _log(f"✅ Imagen convertida a PDF: {ruta_pdf}", None)
            return ruta_pdf
        return None

    except Exception as e:
        _log(f"❌ Error al convertir imagen a PDF: {e}", None)
        return None
### CONVERTIR IMAGEN A PDF
# ══════════════════════════════════════════════════════════════════════════════
# 2. CLIENTES PARA CADA PROVEEDOR
# ══════════════════════════════════════════════════════════════════════════════

class ClienteAPI:
    def __init__(self, nombre: str, modelo: str, api_key: str, cb: Optional[Callable] = None):
        self.nombre = nombre
        self.modelo = modelo
        self.api_key = api_key
        self.cb = cb
        self.historial = []

    def send_message(self, prompt: str, historial: Optional[List[Dict]] = None) -> str:
        raise NotImplementedError

    def send_files(self, rutas: List[str], prompt: str) -> str:
        texto = ""
        for r in rutas:
            if r.lower().endswith('.pdf'):
                txt = extraer_texto_pdf(r)
                if not txt.startswith("[ERROR]"):
                    texto += f"\n--- Contenido de {os.path.basename(r)} ---\n{txt}\n"
        prompt_contenido = f"{prompt}\n\n--- CONTENIDO DE LOS ARCHIVOS ---\n{texto}\n--- FIN DEL CONTENIDO ---"
        return self.send_message(prompt_contenido)

    def _reintentar(self, func, *args, max_retries=3, **kwargs) -> str:
        for intento in range(max_retries):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                error_msg = str(e)
                if any(code in error_msg for code in ["429", "503", "RESOURCE_EXHAUSTED", "UNAVAILABLE", "high demand", "temporarily unavailable"]):
                    match = re.search(r"'retryDelay':\s*'(\d+)s'", error_msg)
                    wait = int(match.group(1)) if match else 8
                    wait = min(wait, 15)
                    _log(f"⏳ {self.nombre}: Error transitorio. Esperando {wait}s (intento {intento+3}/{max_retries})...", self.cb)
                    time.sleep(wait + 1)
                    continue
                else:
                    _log(f"❌ {self.nombre}: {error_msg[:200]}", self.cb)
                    raise
        raise Exception(f"{self.nombre}: Límite de reintentos agotado para error transitorio")


class ClienteGemini(ClienteAPI):
    def __init__(self, api_key: str, modelo: str, cb: Optional[Callable] = None, temperature: float = 0.7, thinking_budget: int = 0):
        super().__init__("Gemini", modelo, api_key, cb)
        self.client = genai.Client(api_key=api_key)
        self.chat = None
        self.archivos_cargados = False
        self.temperature = temperature
        self.thinking_budget = thinking_budget

    def _crear_config_segura(self) -> types.GenerateContentConfig:
        """Genera la configuración de forma segura y compatible con cualquier versión del SDK."""
        config_params = {
            "max_output_tokens": 8192,
            "temperature": self.temperature
        }
        # Si el usuario activó razonamiento, lo inyectamos de forma segura como diccionario crudo
        if self.thinking_budget > 0:
            config_params["thinking_config"] = {"thinking_budget": self.thinking_budget}
            
        return types.GenerateContentConfig(**config_params)

    def iniciar_chat_con_archivos(self, rutas_archivos: List[str]):
        # Intentamos iniciar el chat con la configuración del usuario
        try:
            config = self._crear_config_segura()
            self.chat = self.client.chats.create(model=self.modelo, config=config)
        except Exception as e:
            # Fallback automático: Si el modelo o el SDK antiguo rechaza el thinking_config,
            # conmutamos inmediatamente a una configuración estándar.
            if any(term in str(e).lower() for term in ["thinking", "budget", "unsupported"]):
                _log(f"⚠️ El modelo '{self.modelo}' o tu SDK no admiten razonamiento profundo. Conmutando a modo estándar...", self.cb)
                config = types.GenerateContentConfig(max_output_tokens=8192, temperature=self.temperature)
                self.chat = self.client.chats.create(model=self.modelo, config=config)
            else:
                raise e

        # Bucle original intacto: Lee todos los archivos convertidos a PDF y los envía como bytes binarios
        parts = []
        for r in rutas_archivos:
            if os.path.exists(r):
                datos, mime = leer_bytes(r)
                parts.append(types.Part.from_bytes(data=datos, mime_type=mime))
                    
        parts.append(types.Part.from_text(
            text="Estos son los archivos de la tarea. Los tengo en contexto. "
                 "A partir de ahora, solo recibirás texto; no es necesario que te reenvíe los archivos."
        ))
        self.chat.send_message(parts)
        self.archivos_cargados = True

    def send_message(self, prompt: str, historial: Optional[List[Dict]] = None) -> str:
        if self.chat is None:
            try:
                config = self._crear_config_segura()
                self.chat = self.client.chats.create(model=self.modelo, config=config)
            except Exception as e:
                if any(term in str(e).lower() for term in ["thinking", "budget", "unsupported"]):
                    _log(f"⚠️ El modelo '{self.modelo}' o tu SDK no admiten razonamiento profundo. Conmutando a modo estándar...", self.cb)
                    config = types.GenerateContentConfig(max_output_tokens=8192, temperature=self.temperature)
                    self.chat = self.client.chats.create(model=self.modelo, config=config)
                else:
                    raise e
        return self._reintentar(lambda: self.chat.send_message(prompt).text or "")

    def send_files(self, rutas: List[str], prompt: str) -> str:
        self.iniciar_chat_con_archivos(rutas)
        return "Chat iniciado con archivos"


class ClienteGroq(ClienteAPI):
    def __init__(self, api_key: str, modelo: str, cb: Optional[Callable] = None):
        super().__init__("Groq", modelo, api_key, cb)
        self.client = groq.Client(api_key=api_key)
        self.historial = []

    def send_message(self, prompt: str, historial: Optional[List[Dict]] = None) -> str:
        if historial is not None:
            self.historial = historial
        self.historial.append({"role": "user", "content": prompt})
        respuesta = self._reintentar(
            lambda: self.client.chat.completions.create(
                model=self.modelo,
                messages=self.historial,
                temperature=0.7,
                max_tokens=2048,
            ).choices[0].message.content
        )
        self.historial.append({"role": "assistant", "content": respuesta})
        return respuesta


class ClienteCohere(ClienteAPI):
    def __init__(self, api_key: str, modelo: str, cb: Optional[Callable] = None):
        super().__init__("Cohere", modelo, api_key, cb)
        self.client = cohere.ClientV2(api_key=api_key)
        self.historial = []

    def send_message(self, prompt: str, historial: Optional[List[Dict]] = None) -> str:
        if historial is not None:
            self.historial = historial
        mensajes = self.historial + [{"role": "user", "content": prompt}]
        respuesta = self._reintentar(
            lambda: self.client.chat(
                model=self.modelo,
                messages=mensajes,
                temperature=0.7,
                max_tokens=2048,
            ).message.content[0].text
        )
        self.historial.append({"role": "user", "content": prompt})
        self.historial.append({"role": "assistant", "content": respuesta})
        return respuesta


# ══════════════════════════════════════════════════════════════════════════════
# 3. GESTOR DE PROVEEDORES
# ══════════════════════════════════════════════════════════════════════════════

class ProviderManager:
    def __init__(self, cb: Optional[Callable] = None):
        self.cb = cb
        self.clientes = []
        self.texto_pdf = ""
        self.health = {}

    def registrar_cliente(self, cliente: ClienteAPI):
        self.clientes.append(cliente)
        self.health[cliente.nombre] = {
            "fallos_consecutivos": 0,
            "ultimo_error": None,
            "bloqueado_hasta": 0.0,
            "ultimo_exito": 0.0,
        }

    def _es_error_transitorio(self, error: Exception) -> bool:
        texto = str(error).lower()
        return any(token in texto for token in [
            "429", "503", "resource_exhausted", "unavailable",
            "high demand", "temporarily unavailable", "timeout",
            "rate limit",
        ])

    def _update_health_failure(self, cliente: ClienteAPI, error: Exception):
        estado = self.health.setdefault(cliente.nombre, {
            "fallos_consecutivos": 0,
            "ultimo_error": None,
            "bloqueado_hasta": 0.0,
            "ultimo_exito": 0.0,
        })
        estado["fallos_consecutivos"] += 1
        estado["ultimo_error"] = str(error)
        if self._es_error_transitorio(error):
            backoff = min(60, 5 * estado["fallos_consecutivos"])
            estado["bloqueado_hasta"] = time.time() + backoff
            _log(f"⏳ {cliente.nombre} se bloquea temporalmente por {backoff}s tras error transitorio.", self.cb)
        else:
            estado["bloqueado_hasta"] = time.time() + 5
        return estado

    def _update_health_success(self, cliente: ClienteAPI):
        estado = self.health.setdefault(cliente.nombre, {
            "fallos_consecutivos": 0,
            "ultimo_error": None,
            "bloqueado_hasta": 0.0,
            "ultimo_exito": 0.0,
        })
        estado["fallos_consecutivos"] = 0
        estado["ultimo_error"] = None
        estado["bloqueado_hasta"] = 0.0
        estado["ultimo_exito"] = time.time()

    def _ordenar_clientes_para_intentos(self) -> List[ClienteAPI]:
        ahora = time.time()
        activos = []
        for cliente in self.clientes:
            # EXCLUIR al cliente de respaldo del flujo e intentos estándar
            if getattr(cliente, "es_respaldo", False):
                continue
            estado = self.health.get(cliente.nombre, {})
            bloqueado = estado.get("bloqueado_hasta", 0) > ahora
            activos.append((bloqueado, estado.get("fallos_consecutivos", 0), -(estado.get("ultimo_exito", 0)), cliente))
        activos.sort(key=lambda item: (item[0], item[1], item[2]))
        return [item[3] for item in activos]
    
    def _puede_intentar(self, cliente: ClienteAPI) -> bool:
        estado = self.health.get(cliente.nombre, {})
        return estado.get("bloqueado_hasta", 0) <= time.time()

    def send_message(self, prompt: str, historial: Optional[List[Dict]] = None) -> tuple[str, str]:
        clientes_ordenados = self._ordenar_clientes_para_intentos()
        ultimo_error = None

        for cliente in clientes_ordenados:
            if not self._puede_intentar(cliente):
                _log(f"⏳ Saltando {cliente.nombre} mientras se recupera de errores anteriores.", self.cb)
                continue
            try:
                _log(f"🔄 Intentando con {cliente.nombre}...", self.cb)

                if isinstance(cliente, ClienteGemini):
                    prompt_final = prompt
                else:
                    if self.texto_pdf:
                        max_chars = 3500
                        texto_limitado = self.texto_pdf[:max_chars]
                        if len(self.texto_pdf) > max_chars:
                            texto_limitado += "\n... [texto truncado para ajustarse a límites de Groq/Cohere]"
                        prompt_final = f"{prompt}\n\n--- CONTENIDO DE LOS PDFS (RESUMIDO) ---\n{texto_limitado}\n--- FIN DEL CONTENIDO ---"
                    else:
                        prompt_final = prompt

                respuesta = cliente.send_message(prompt_final, historial)
                self._update_health_success(cliente)
                _log(f"✅ {cliente.nombre} respondió correctamente.", self.cb)
                return respuesta, cliente.nombre
            except Exception as e:
                ultimo_error = e
                _log(f"⚠️ {cliente.nombre} falló: {str(e)[:200]}", self.cb)
                self._update_health_failure(cliente, e)
                continue

        if ultimo_error:
            raise Exception(f"Todos los proveedores fallaron: {ultimo_error}")
        raise Exception("No hay proveedores disponibles para procesar la solicitud.")

    def send_files(self, rutas: List[str], prompt: str) -> tuple[str, str]:
        # Guardamos la lista de rutas para que el failover sepa qué archivos anexar
        self.rutas_originales_guardadas = rutas

        texto_completo = ""
        for r in rutas:
            if r.lower().endswith('.pdf'):
                txt = extraer_texto_pdf(r)
                if not txt.startswith("[ERROR]"):
                    texto_completo += f"\n--- Archivo: {os.path.basename(r)} ---\n{txt}\n"
        self.texto_pdf = texto_completo

        for cliente in self.clientes:
            if isinstance(cliente, ClienteGemini):
                try:
                    _log(f"🔄 Iniciando chat de Gemini con archivos...", self.cb)
                    cliente.iniciar_chat_con_archivos(rutas)
                    self._update_health_success(cliente)
                    _log(f"✅ Gemini chat iniciado con archivos.", self.cb)
                except Exception as e:
                    _log(f"⚠️ Gemini falló al iniciar chat: {str(e)[:200]}", self.cb)
                    self._update_health_failure(cliente, e)

        return "", "files_loaded"


# ══════════════════════════════════════════════════════════════════════════════
# FUNCIONES AUXILIARES DE CREACIÓN DE CLIENTES
# ══════════════════════════════════════════════════════════════════════════════

def crear_cliente_gemini(api_key: str, modelo: str, cb: Optional[Callable] = None, 
    temperature: float = 0.7, thinking_budget: int = 0) -> ClienteGemini:
    return ClienteGemini(api_key, modelo, cb, temperature, thinking_budget)

def crear_cliente_groq(api_key: str, modelo: str, cb: Optional[Callable] = None) -> ClienteGroq:
    return ClienteGroq(api_key, modelo, cb)

def crear_cliente_cohere(api_key: str, modelo: str, cb: Optional[Callable] = None) -> ClienteCohere:
    return ClienteCohere(api_key, modelo, cb)


# ══════════════════════════════════════════════════════════════════════════════
# CLASIFICACIÓN
# ══════════════════════════════════════════════════════════════════════════════

_PROMPT_CLASIF = """\
Basado en los archivos que ya tienes en contexto Y en las instrucciones del usuario,
responde EXACTAMENTE con dos líneas (nada más):
curso: [ingles|matematicas|general]
longitud: [LARGO|CORTO]

Instrucciones del usuario (esto tiene PRIORIDAD sobre el formato del archivo):
{instrucciones}

Criterios de clasificación DETALLADOS:

- ingles: Solo si la tarea es EXCLUSIVAMENTE de idioma inglés (gramática, emails, diálogos en inglés;etc).
- matematicas: Solo si la tarea pide explícitamente resolver ejercicios numéricos de matemática, estadística, cálculo, álgebra, probabilidad o análisis de datos con tablas numéricas para resolver.
OJO: si aparte de pedir calculos numéricos, la tarea también pide redactar contenido formal (indice, introduccion,marco teorico, metodologia, conclusion, referencias;etc), se clasifica como "general" (no matemáticas) porque el foco no es solo resolver ejercicios numéricos, sino también interpretar y redactar.
- general: CUALQUIER OTRA COSA por ejemplo: Ciudadanía, Administración, Economía, Historia, Filosofía, Psicología, Sociología, Derecho, Comunicación, Proyectos, Tesis cualitativas, etc.

- LARGO: si la consigna pide 3 o más páginas, o estructura de informe con secciones formales (Introducción, Metodología, Resultados, Conclusión; etc), o es un proyecto/tarea que claramente requiere un desarrollo extenso.
- CORTO: si es tarea de 1-2 páginas, o son ejercicios puntuales con cálculo numérico (incluso
  si tiene varias etapas/problemas, cada uno con datos para calcular).

IMPORTANTE: Si el usuario pide resolver UN problema/ejercicio específico (ej: "resuelve el
Problema 1"), clasifica SOLO en base a ese problema específico (QUE POSIBLEMENTE TENDRÁ SUBPROBLEMAS ADENTRO), ignorando el resto del archivo.
"""

def clasificar(manager: ProviderManager, instrucciones: str = "", cb=None) -> tuple[str, str]:
    _log("🔍 Clasificando tipo y longitud...", cb)
    prompt_final = _PROMPT_CLASIF.format(
        instrucciones=instrucciones.strip() if instrucciones and instrucciones.strip()
        else "(El usuario no especificó instrucciones adicionales; clasifica según el archivo completo.)"
    )
    try:
        resp, _ = manager.send_message(prompt_final)
    except Exception as e:
        _log(f"⚠️ Error en clasificación: {e}", cb)
        return "general", "CORTO"
    _log(f"📋 Clasificación: {resp}", cb)
    curso = "general"
    longitud = "CORTO"
    for linea in resp.lower().split("\n"):
        linea = linea.strip()
        if linea.startswith("curso:"):
            val = linea.split(":", 1)[1].strip()
            if val in ("ingles", "matematicas", "general"):
                curso = val
        elif linea.startswith("longitud:"):
            val = linea.split(":", 1)[1].strip()
            if val in ("largo", "corto"):
                longitud = val.upper()

    # ── RED DE SEGURIDAD: si las instrucciones mencionan palabras de cálculo
    # estadístico/matemático explícito, forzamos 'matematicas' aunque la IA
    # haya clasificado distinto. Esto evita que informes con etapas numeradas
    # (1.1, 1.2, 1.3) sean tratados como "general" solo por su formato.
    señales_calculo = [
        "resuelve", "calcula", "chi-cuadrado", "chi cuadrado", "mann-whitney",
        "mann whitney", "kruskal", "wallis", "estadístico", "estadistico",
        "hipótesis nula", "hipotesis nula", "valor crítico", "valor critico",
        "anova", "regresión", "regresion", "intervalo de confianza",
        "prueba de hipótesis", "prueba de hipotesis", "distribución", "distribucion",
    ]
    texto_check = instrucciones.lower() if instrucciones else ""
    if curso != "matematicas" and any(s in texto_check for s in señales_calculo):
        _log(f"🔧 Corrección automática: instrucciones mencionan cálculo estadístico → forzando curso=matematicas", cb)
        curso = "matematicas"

    _log(f"✅ Curso: {curso} | Longitud: {longitud}", cb)
    return curso, longitud


# ══════════════════════════════════════════════════════════════════════════════
# FLUJO MATEMÁTICAS
# ══════════════════════════════════════════════════════════════════════════════

_PROMPT_CONTAR_EJERCICIOS = """\
Basado en los archivos que ya tienes en contexto.

Instrucciones del usuario: {instrucciones}

Si el usuario pidió resolver UN problema O LO MENCIONÓ EXPLÍCITAMENTE (ej: "resuelve el Problema 1"), cuenta SOLO las SUB-ETAPAS o SUB-EJERCICIOS que tiene ESE problema específico (ignora los demás problemas del archivo).
responde SOLO con el número de SUB-EJERCICIOS o ETAPAS que tiene ESE problema específico
(ignora los demás problemas del archivo). POR EJEMPLO SI EL PROBLEMA 1 TIENE 3 ETAPAS CON 3 EJERCICIOS CADA UNA, DEBERAS RESPONDER: 9

Si el usuario no especificó un problema concreto, cuenta TODOS los ejercicios del archivo.

Responde SOLO con un número entero. Por ejemplo: 8
No escribas nada más.
"""

_PROMPT_RESOLVER_EJERCICIO = """\
Eres un experto en matemáticas, estadística y cálculo.

Instrucciones del usuario: {instrucciones}

REGLAS CRÍTICAS PARA ECUACIONES (EVITA ERRORES DE PARSEO EN WORD):
1. FORMATO ÚNICO [EQ: ...]: Usa exclusivamente el formato [EQ: código_LaTeX] para fórmulas complejas, fracciones, sistemas de ecuaciones enteros, límites o matrices. NO uses signos de dólar ($ o $$).
2. PROHIBIDO USAR CORCHETES DENTRO DE [EQ: ...]: Está estrictamente prohibido usar corchetes cuadrados `[` o `]` dentro del código LaTeX (esto rompe el parser del documento). 
   - Para representar matrices o vectores, usa obligatoriamente el entorno pmatrix con doble escape: `\\\\begin{{pmatrix}} a & b \\\\\\\\ c & d \\\\end{{pmatrix}}`. NUNCA uses bmatrix, ni corchetes manuales.
   - Si necesitas agrupar términos en una fórmula, usa paréntesis normales o llaves de LaTeX del tipo `\\\\left( ... \\\\right)` o `\\\\left\\{{ ... \\\\right\\}}`.
3. VARIABLES Y ASIGNACIONES SIMPLES EN TEXTO NORMAL: NO uses el formato [EQ: ...] para variables, coordenadas o resultados simples independientes (ej: x = 3, y = 2, z = 1 o x = 83/119 deben escribirse en texto plano normal (salvo que sean elevados al cuadrado, cubo; etc), sin envolverlos en [EQ: ...]). Solo envuelve ecuaciones largas con operaciones complejas(ejemplo: elevados al cuadrado, cubo, elementos que necesitan estar en formato ecuacion para visualización correcta; etc).
4. EVITA CORCHETES EN NOMBRES DE FILAS/COFACTORES: Si usas variables con subíndices o indicadores de operación como F2 <- F2 - 4F1 o C11, escríbelas siempre como texto plano normal fuera de bloques [EQ: ...].
5. AGRUPACIÓN DE RESULTADOS: Al dar la respuesta o resultado final de un sistema, agrúpalos en una sola línea de texto normal clara (ej. "Resultado: x = 3, y = 2, z = 1") en lugar de crear múltiples bloques de ecuaciones individuales pegados.

REGLAS PARA EL CONTENIDO — OBLIGATORIO:
- Si el usuario pidió o mencionó una parte o problema específico, resuelve ese problema completo (preguntas respondidas, calculos realizados, sub ejercicios resueltos;etc) (ignora el resto del archivoque no corresponde a lo pedido).
- Cada Ejercicio que resuelvas DEBE ser ubicable (pero debes adaptarte al contexto por ejemplo: "Etapa 1 - Ejercicio 1", "Fase 1 - Ejercicio 1", "Problema 1 - Subejercicio 1", etc) para que el usuario pueda ubicarlo en el archivo original.
- Copia el enunciado (pero no empieces por "Enunciado: ......", Mejor redacta directamente su contenido) y a que fase o etapa pertenece (si es que aplica), plantea hipótesis (si es que aplica), CALCULA PASO A PASO CON LOS NÚMEROS REALES
  (no te quedes solo en la fórmula teórica — debes EJECUTAR el cálculo numérico completo:
  reemplaza los valores, opera, obtén el resultado numérico final).
- Si la prueba requiere asignar rangos (Mann-Whitney, Kruskal-Wallis), DEBES escribir la
  tabla de rangos REAL con cada dato ordenado y su rango asignado — no la omitas.
- Muestra el valor numérico final del estadístico calculado (χ², U, H, etc.) — un número concreto,
  no una fórmula sin resolver.
- Da la decisión estadística (rechaza/no rechaza H0) comparando el número calculado contra el crítico.
- NO dejes notas ni comentarios sueltos.
- Usa los datos EXACTOS del archivo (los tienes en contexto).
"""

def contar_ejercicios(manager: ProviderManager, instrucciones: str = "", cb=None) -> int:
    _log("🔢 Contando ejercicios...", cb)
    prompt = _PROMPT_CONTAR_EJERCICIOS.format(
        instrucciones=instrucciones.strip() if instrucciones and instrucciones.strip()
        else "(sin instrucciones específicas; cuenta todos los ejercicios del archivo)"
    )
    try:
        resp, _ = manager.send_message(prompt)
    except Exception:
        return 1
    m = re.search(r'\d+', resp)
    n = int(m.group()) if m else 1
    n = max(1, min(n, 20))
    _log(f"✅ Ejercicios detectados: {n}", cb)
    return n

def resolver_ejercicio(manager: ProviderManager, numero: int, total: int, instrucciones: str = "", cb=None) -> str:
    _log(f"✍️  Resolviendo ejercicio {numero}/{total}...", cb)
    prompt_base = _PROMPT_RESOLVER_EJERCICIO.format(
        instrucciones=instrucciones.strip() if instrucciones and instrucciones.strip()
        else "(sin instrucciones específicas; resuelve según el orden del archivo)"
    )
    prompt = (
        f"{prompt_base}\n\n"
        f"TAREA: Resuelve ÚNICAMENTE el EJERCICIO/ETAPA NÚMERO {numero}.\n"
        f"(Total: {total})\n"
        f"Escribe el título: '## Ejercicio {numero}' al inicio.\n"
    )
    # Propagar el error al bucle superior
    resultado, _ = manager.send_message(prompt)
    return resultado

def generar_matematicas(manager: ProviderManager, instrucciones, cb=None, stop_check=None) -> str:
    _log("📐 Modo MATEMÁTICAS activado", cb)
    n = contar_ejercicios(manager, instrucciones, cb)
    partes = []
    
    i = 1
    while i <= n:
        if stop_check:
            stop_check()
        try:
            sol = resolver_ejercicio(manager, i, n, instrucciones, cb)
            partes.append(sol)
            _log(f"   ✅ Ejercicio {i}/{n} completado ({len(sol)} chars)", cb)
            i += 1
        except Exception as e:
            _log(f"⚠️ Error crítico en ejercicio matemático {i}: {e}", cb)
            
            cliente_respaldo = next((c for c in manager.clientes if getattr(c, "es_respaldo", False)), None)
            if cliente_respaldo:
                _log("🚨 ¡Clave API Principal agotada! Activando API de Respaldo de Gemini para Matemáticas...", cb)
                _log("⏳ Pausando el flujo durante 5 segundos para estabilizar la cuota...", cb)
                time.sleep(5)
                
                try:
                    # 1. Guardar lo avanzado en matemáticas
                    texto_incompleto = "\n\n".join(partes)
                    ruta_incompleta_docx = "temp_borrador_incompleto.docx"
                    guardar_word(texto_incompleto, "matematicas", ruta_incompleta_docx)
                    
                    # 2. Convertir a PDF
                    _log("📄 Convirtiendo borrador numérico a PDF...", cb)
                    ruta_incompleta_pdf = _word_a_pdf(ruta_incompleta_docx, cb)
                    
                    # 3. Anexar borrador
                    rutas_con_respaldo = []
                    originales = getattr(manager, "rutas_originales_guardadas", [])
                    rutas_con_respaldo.extend(originales)
                    
                    if ruta_incompleta_pdf and os.path.exists(ruta_incompleta_pdf):
                        rutas_con_respaldo.append(ruta_incompleta_pdf)
                        _log("🔗 Borrador matemático acoplado al portafolio de archivos del proveedor de respaldo.", cb)
                    
                    # 4. Inicializar respaldo
                    _log("🔄 Inicializando chat de contingencia con API de Respaldo...", cb)
                    cliente_respaldo.iniciar_chat_con_archivos(rutas_con_respaldo)
                    
                    # 5. Intercambiar clientes
                    # IMPORTANTE: igual que en generar_general(), se quita la marca 'es_respaldo'
                    # porque _ordenar_clientes_para_intentos() excluye a los clientes de respaldo.
                    # Si se deja la marca, al ser ahora el único cliente registrado, cada intento
                    # devuelve una lista vacía -> "No hay proveedores disponibles" -> loop infinito.
                    cliente_respaldo.es_respaldo = False
                    manager.clientes = [cliente_respaldo]
                    manager.health = {cliente_respaldo.nombre: {
                        "fallos_consecutivos": 0, "ultimo_error": None, "bloqueado_hasta": 0.0, "ultimo_exito": time.time()
                    }}
                    
                    try: os.remove(ruta_incompleta_docx)
                    except: pass
                    try: os.remove(ruta_incompleta_pdf)
                    except: pass
                    
                    _log("🚀 API de Respaldo para Matemáticas lista. Reanudando...", cb)
                    continue
                except Exception as ex_respaldo:
                    _log(f"❌ Falló el despliegue del proveedor de respaldo matemático: {ex_respaldo}", cb)
                    break
            else:
                _log("❌ No se configuró ninguna clave de respaldo (API 2). Abortando matemáticas.", cb)
                break
                
    return "\n\n".join(partes)


# ══════════════════════════════════════════════════════════════════════════════
# EXTRACCIÓN DE DATOS POR SECCIÓN
# ══════════════════════════════════════════════════════════════════════════════

def extraer_datos_por_seccion(manager: ProviderManager, seccion, cb=None):
    prompt_extract = f"""
    Eres un analista de datos. Extrae toda la información técnica, numérica, tablas y 
    cifras relevantes del PDF para la sección: '{seccion}'.
    
    REGLAS:
    1. Si hay tablas, represéntalas en formato Markdown.
    2. Si no hay datos suficientes, escribe 'DATOS_NO_DISPONIBLES'.
    3. Si tienes conocimiento externo que respalde esto, agrégalo brevemente.
    4. NO redactes párrafos, solo entrega datos, tablas y hechos.
    """
    try:
        resp, _ = manager.send_message(prompt_extract)
        return resp
    except Exception:
        return "DATOS_NO_DISPONIBLES"


# ══════════════════════════════════════════════════════════════════════════════
# FLUJO GENERAL / INGLÉS
# ══════════════════════════════════════════════════════════════════════════════

_PROMPT_PLAN = """\
Eres un experto en estructuración de tareas, proyectos, tesis e informes académicos. Analiza el archivo de la consigna y la rúbrica.

Tu tarea es inferir la estructura de secciones que el docente espera ver para obtener la máxima calificación.

REGLAS OBLIGATORIAS:
- Devuelve UNA LISTA PLANA numerada con TODOS los niveles jerárquicos.
- Incluye las subsecciones como elementos separados (ej: 1., 1.1, 1.2, 2., 2.1, 2.2, etc.).
- NO agrupes el contenido de las subsecciones dentro de la sección principal. Cada subsección debe ser un elemento independiente en la lista.
- Si la rúbrica menciona varios temas dentro de una sección, crea una subsección para cada tema.
- Devuelve SOLO la lista numerada (con subsecciones si aplica). Nada más.
- Las subsecciones que necesiten desarrollo detallado (más de 2 páginas) deben tener numeracion estilo "2.1." (su numero tendrá un punto al final) y las que requieran desarrollo inferior a 2 paginas estilo "2.1" (su numero no tendrá un punto al final).
- Si la consigna no pide una estructura especifica, debes crear una apartir de lo existente en la tarea:
 Por Ejemplo:
1. Problema 1:
1.1. Etapa 1
1.2. Etapa 2
1.3. Etapa 3
1.4 interpretacion y conclusión
2. Problema 2:
2.1. Etapa 1
2.2. Etapa 2
2.3. Etapa 3
2.4 interpretacion y conclusión
- Salvo que la tarea o rubrica exija una estructura mas formal:
EJEMPLO CORRECTO (pero ajustable según la tarea):
1. Introducción
1.1 Contexto
1.2 Objetivos
1.3 Justificación
1.4 Conclusion conclusiva de la introducción
2. Marco Teórico
2.1. Teoría de la Dependencia
2.2. Teoría del Desarrollo
2.3. Estado del Arte
3. Metodología
3.1 Diseño de investigación
3.2 Población y muestra
3.3 Técnicas de recolección
4. Resultados
4.1. Análisis descriptivo
4.2. Pruebas de hipótesis
5. Discusión
6. Conclusiones
7. Recomendaciones
8. Referencias

Ahora, genera la estructura específica para este caso.
"""

def planificar_secciones(manager: ProviderManager, cb=None) -> list[str]:
    _log("📐 Planificando secciones (incluyendo subsecciones)...", cb)
    try:
        resp, _ = manager.send_message(_PROMPT_PLAN)
        if resp.startswith("ERROR"):
            raise Exception(resp)
        _log(f"📋 Plan:\n{resp}", cb)
        secciones = []
        for linea in resp.split("\n"):
            linea = linea.strip()
             # Captura 1., 1.1, 1.1.1, 2., 2.1, etc. "(r'^(\d+(?:\.\d+)*)[\.\)]?\s+(.+)$', linea)" <--el simbolo ? hace la diferencia
            m = re.match(r'^(\d+(?:\.\d+)*)[\.\)]\s+(.+)$', linea)
            if m:
                numero_completo = m.group(1)
                titulo = m.group(2).strip()
                secciones.append(f"{numero_completo}. {titulo}")   # <-- Conserva el número
        if not secciones:
            secciones = ["Introducción", "Desarrollo", "Análisis", "Conclusiones", "Referencias"]
        _log(f"✅ {len(secciones)} secciones/subsecciones detectadas.", cb)
        return secciones
    except Exception:
        _log("⚠️ Usando secciones genéricas.", cb)
        return ["Introducción", "Desarrollo", "Análisis", "Conclusiones", "Referencias"]

_REGLAS_REDACCION = """\
REGLAS OBLIGATORIAS:
- Redacta SOLO el contenido de la sección indicada. Nada más.
- NO escribas el título de la sección (el programa lo añade).
- NO dejes notas ni comentarios.
- NO hagas preguntas al usuario.
- Incluye DATOS REALES (prioriza usar los datos brindados por los archivos en contexto) y cifras.
- NO incluyas tablas en formato markdown o gráficos descriptivos a menos que la rúbrica, la consigna o las instrucciones lo indiquen explícitamente, o que esta sección lo amerite de forma urgente (por ejemplo: para organizar datos numéricos complejos, reportes financieros, rendimientos comparativos;etc).
- NO añadas secciones, conceptos o análisis que no se pidan.
- NO omitas secciones, conceptos o análisis que se pidan.
- Cada sección debe tener contenido valioso que responda y cumpla su funcion con lo que pide (la rúbrica o consigna) para esa sección.
- Por ejemplo, si la rubrica o tarea mencionó calcular X variables o responder Y preguntas, asegúrate de que la sección incluya esos cálculos o respuestas específicas.
- Tono académico, formal. Usa conectores lógicos.
- Si la sección incluye tablas, represéntalas en formato markdown.
"""

def redactar_seccion(manager: ProviderManager, seccion, numero, total, contexto_ant, prompt_base, cb=None) -> str:
    _log(f"✍️  Redactando sección {numero}/{total}: '{seccion}'", cb)
    datos_seccion = extraer_datos_por_seccion(manager, seccion, cb)
    contexto_str = ""
    if contexto_ant.strip():
        contexto_str = f"\n\nCONTEXTO previo (NO repetir):\n[...]\n{contexto_ant[-1500:]}\n[fin]"
    prompt = (
        f"{prompt_base}\n\n"
        f"DATOS TÉCNICOS EXTRAÍDOS PARA ESTA SECCIÓN:\n{datos_seccion}\n\n"
        f"{_REGLAS_REDACCION}\n"
        f"{contexto_str}\n\n"
        f"TAREA: Redacta la sección «{seccion}» (sección {numero} de {total}). "
    )
# try:  --- Retiramos el try-except interno para delegar la conmutación al flujo superior
    resultado, _ = manager.send_message(prompt)
    return resultado
    #except Exception as e:
    #   if any(code in str(e) for code in ["429", "503", "RESOURCE_EXHAUSTED"]):
    #       return f"[Sección '{seccion}' no generada por saturación/cuota.]"
    #    return f"[Error: {str(e)[:100]}]"

def generar_general(manager: ProviderManager, instrucciones, prompt_base, longitud, cb=None, stop_check=None) -> str:
    if longitud == "CORTO":
        _log("📋 Tarea CORTA — una llamada...", cb)
        datos_generales = extraer_datos_por_seccion(manager, "todo el trabajo", cb)
        prompt = (
            f"{prompt_base}\n\n"
            f"DATOS EXTRAÍDOS DEL PDF:\n{datos_generales}\n\n"
            f"{_REGLAS_REDACCION}\n\n"
            "Redacta el trabajo completo ahora."
        )
        resultado, _ = manager.send_message(prompt)
        return resultado
    else:
        _log("📚 Tarea LARGA — sección por sección...", cb)
        secciones = planificar_secciones(manager, cb)
        partes_texto = []
        contexto = ""
        total = len(secciones)
        
        i = 1
        while i <= total:
            if stop_check:
                stop_check()
            sec = secciones[i-1]
            try:
                # Intento normal con el proveedor primario
                contenido = redactar_seccion(manager, sec, i, total, contexto, prompt_base, cb)
                bloque = f"# {sec}\n{contenido}"
                partes_texto.append(bloque)
                contexto = "\n\n".join(partes_texto)
                _log(f"   ✅ {i}/{total} completada ({len(contenido)} chars)", cb)
                i += 1
            except Exception as e:
                _log(f"⚠️ Error crítico al generar sección '{sec}': {e}", cb)
                
                # Buscar si hay un cliente de respaldo registrado
                cliente_respaldo = next((c for c in manager.clientes if getattr(c, "es_respaldo", False)), None)
                
                if cliente_respaldo:
                    _log("🚨 ¡Clave API Principal agotada! Activando API de Respaldo de Gemini...", cb)
                    _log("⏳ Pausando el flujo durante 5 segundos para estabilizar la cuota...", cb)
                    time.sleep(5)
                    
                    try:
                        # 1. Guardar lo que llevamos redactado hasta el momento en un Word temporal
                        texto_incompleto = "\n\n".join(partes_texto)
                        ruta_incompleta_docx = "temp_borrador_incompleto.docx"
                        guardar_word(texto_incompleto, "general", ruta_incompleta_docx)
                        
                        # 2. Convertir el borrador incompleto a PDF
                        _log("📄 Convirtiendo borrador redactado a PDF para análisis de contexto...", cb)
                        ruta_incompleta_pdf = _word_a_pdf(ruta_incompleta_docx, cb)
                        
                        # 3. Anexar el borrador a la lista de archivos originales
                        rutas_con_respaldo = []
                        originales = getattr(manager, "rutas_originales_guardadas", [])
                        rutas_con_respaldo.extend(originales)
                        
                        if ruta_incompleta_pdf and os.path.exists(ruta_incompleta_pdf):
                            rutas_con_respaldo.append(ruta_incompleta_pdf)
                            _log("🔗 Borrador incompleto acoplado al portafolio de archivos del proveedor de respaldo.", cb)
                        
                        # 4. Inicializar el contexto del cliente de respaldo con la costura completa
                        _log("🔄 Inicializando chat de contingencia con API de Respaldo...", cb)
                        cliente_respaldo.iniciar_chat_con_archivos(rutas_con_respaldo)
                        
                        # 5. Sustituir la lista de proveedores del Manager para usar únicamente el respaldo
                        # IMPORTANTE: se quita la marca 'es_respaldo' porque _ordenar_clientes_para_intentos()
                        # excluye a cualquier cliente marcado como respaldo. Si se deja la marca, al ser ahora
                        # el ÚNICO cliente, la lista de candidatos queda vacía en cada intento -> "No hay
                        # proveedores disponibles" -> se vuelve a activar "el respaldo" -> loop infinito.
                        cliente_respaldo.es_respaldo = False
                        manager.clientes = [cliente_respaldo]
                        manager.health = {cliente_respaldo.nombre: {
                            "fallos_consecutivos": 0, "ultimo_error": None, "bloqueado_hasta": 0.0, "ultimo_exito": time.time()
                        }}
                        
                        # Limpieza física de los borradores temporales
                        try: os.remove(ruta_incompleta_docx)
                        except: pass
                        try: os.remove(ruta_incompleta_pdf)
                        except: pass
                        
                        _log("🚀 Proveedor de Respaldo sincronizado correctamente. Reanudando redacción...", cb)
                        # No avanzamos 'i' para que reintente la generación de la sección fallida con la nueva API
                        continue
                    except Exception as ex_respaldo:
                        _log(f"❌ Falló el despliegue del proveedor de respaldo: {ex_respaldo}", cb)
                        break
                else:
                    _log("❌ No se configuró ninguna clave de respaldo (API 2). Abortando proceso.", cb)
                    break
        
        return "\n\n".join(partes_texto)


# ══════════════════════════════════════════════════════════════════════════════
# LIMPIEZA Y GUARDADO
# ══════════════════════════════════════════════════════════════════════════════

def limpiar(texto: str) -> str:
    for patron in [r'\[CONTENIDO_LARGO\]', r'\[CONTENIDO_CORTO\]',
                   r'\[PARTE_FIN\]', r'\[FIN\]']:
        texto = re.sub(patron, '', texto, flags=re.IGNORECASE)
    texto = re.sub(r'\n{3,}', '\n\n', texto)
    texto = re.sub(r' +', ' ', texto)
    return texto.strip()

def guardar_word(texto: str, tipo: str, ruta: str = "Tarea_Realizada.docx") -> str:
    tipo_generador = "estadistica" if tipo == "matematicas" else tipo
    texto_limpio = limpiar(texto)
    if tipo_generador == "estadistica":
        gen = GeneradorEstadistica(nombre_archivo=ruta)
        gen.procesar_todo(texto_limpio)
        gen.guardar()
    else:
        gen = GeneradorWord(nombre_archivo=ruta, tipo_tarea=tipo_generador)
        gen.agregar_texto(texto_limpio)
        gen.guardar()
    if platform.system() == "Windows":
        optimizar_word_con_win32com(ruta)
    return ruta

def optimizar_word_con_win32com(ruta_word: str, cb=None) -> bool:
    """
    Optimiza el Word de forma multiplataforma.
    Limpia saltos de línea repetidos de forma nativa en cualquier SO,
    y ejecuta OMaths.BuildUp() únicamente si corre sobre Windows.
    """
    import os
    import platform
    from docx import Document

    if not os.path.exists(ruta_word):
        return False

    # 1. Limpieza nativa multiplataforma (reemplazo seguro de "^p^p")
    try:
        doc = Document(ruta_word)
        p_eliminar = []
        vacio_previo = False
        
        for p in doc.paragraphs:
            if not p.text.strip():  # Párrafo vacío
                if vacio_previo:
                    p_eliminar.append(p)
                else:
                    vacio_previo = True
            else:
                vacio_previo = False

        for p in p_eliminar:
            p_element = p._element
            p_element.getparent().remove(p_element)
            
        doc.save(ruta_word)
    except Exception as e:
        _log(f"⚠️ Error en limpieza nativa de párrafos: {e}", cb)

    # 2. Compilación matemática exclusiva de Windows (win32com)
    if platform.system() == "Windows":
        try:
            import pythoncom
            import win32com.client
            
            pythoncom.CoInitialize()
            word_app = win32com.client.Dispatch("Word.Application")
            word_app.Visible = False
            doc_win = word_app.Documents.Open(os.path.abspath(ruta_word))
            doc_win.OMaths.BuildUp()
            doc_win.Save()
            doc_win.Close()
            word_app.Quit()
            pythoncom.CoUninitialize()
            return True
        except Exception as e:
            _log(f"⚠️ OMaths BuildUp de Windows omitido: {e}", cb)
            
    return True



# ══════════════════════════════════════════════════════════════════════════════
# EVALUACIÓN POR IA
# ══════════════════════════════════════════════════════════════════════════════

_PROMPT_EVALUACION = """\
Eres un docente ESTRICTO y EXIGENTE.

Acabas de recibir el PDF con el trabajo del alumno.
Ya tienes en contexto la consigna y la rúbrica originales del curso.

Evalúa el trabajo del alumno comparándolo con la rúbrica.
Sé estricto: penaliza lo que falta, valora lo que está bien hecho.

Responde EXACTAMENTE con este formato (sin agregar nada antes ni después):

NOTA: [número del 0 al 20]
JUSTIFICACIÓN:
[Escribe aquí tu evaluación detallada: qué cumplió, qué le faltó, por qué esa nota.]
"""

def _word_a_pdf(ruta_word: str, cb=None) -> str | None:
    _log("📄 Convirtiendo Word a PDF...", cb)
    directorio = os.path.dirname(os.path.abspath(ruta_word))
    return convertir_docx_a_pdf(ruta_word, directorio)

def evaluar_tarea_con_ia(manager: ProviderManager, modelo_ia: str, ruta_word: str, cb=None) -> dict:
    _log("🎓 Iniciando evaluación docente IA...", cb)
    ruta_pdf = _word_a_pdf(ruta_word, cb)
    if ruta_pdf is None:
        return {"nota": "?", "justificacion": "No se pudo generar PDF.", "texto_completo": ""}
    try:
        lector = PdfReader(ruta_pdf)
        texto_pdf = "\n".join(p.extract_text() for p in lector.pages if p.extract_text())
        prompt = f"{_PROMPT_EVALUACION}\n\n--- TRABAJO ---\n{texto_pdf}\n--- FIN ---"
        resp, _ = manager.send_message(prompt)
    except Exception as e:
        _log(f"⚠️ Error en evaluación: {e}", cb)
        return {"nota": "?", "justificacion": f"Error: {e}", "texto_completo": ""}
    finally:
        try:
            if ruta_pdf and ruta_pdf != ruta_word.replace(".docx", ".pdf"):
                os.remove(ruta_pdf)
        except:
            pass

    nota = "?"
    justificacion = resp
    m_nota = re.search(r'NOTA:\s*(\d+(?:\.\d+)?)', resp, re.IGNORECASE)
    if m_nota:
        nota = m_nota.group(1)
    m_just = re.search(r'JUSTIFICACI[OÓ]N:\s*(.*)', resp, re.IGNORECASE | re.DOTALL)
    if m_just:
        justificacion = m_just.group(1).strip()
    _log(f"✅ Evaluación — Nota: {nota}/20", cb)
    return {"nota": nota, "justificacion": justificacion, "texto_completo": resp}


# ══════════════════════════════════════════════════════════════════════════════
# GENERADOR EXCEL
# ══════════════════════════════════════════════════════════════════════════════

def generar_excel_desde_texto(texto: str, ruta_salida: str = "Tarea_Realizada_datos.xlsx") -> str:
    lineas = texto.split('\n')
    tablas = []
    tabla_actual = []
    dentro_tabla = False

    for linea in lineas:
        if '|' in linea and linea.strip().startswith('|') and linea.strip().endswith('|'):
            dentro_tabla = True
            tabla_actual.append(linea.strip())
        elif dentro_tabla:
            if tabla_actual:
                tablas.append(tabla_actual)
                tabla_actual = []
            dentro_tabla = False
    if tabla_actual:
        tablas.append(tabla_actual)

    if not tablas:
        wb = Workbook()
        ws = wb.active
        ws.title = "Sin_tablas"
        ws['A1'] = "No se encontraron tablas en el texto."
        wb.save(ruta_salida)
        return ruta_salida

    header_fill = PatternFill(start_color="2E74B5", end_color="2E74B5", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    border = Border(
        left=Side(style='thin', color='000000'),
        right=Side(style='thin', color='000000'),
        top=Side(style='thin', color='000000'),
        bottom=Side(style='thin', color='000000')
    )
    even_fill = PatternFill(start_color="EBF3FB", end_color="EBF3FB", fill_type="solid")

    with pd.ExcelWriter(ruta_salida, engine='openpyxl') as writer:
        validacion_data = []

        for idx, bloque in enumerate(tablas, start=1):
            lineas_limpias = []
            for line in bloque:
                if re.match(r'^[\|\s\-:]+$', line):
                    continue
                partes = [p.strip() for p in line.strip('|').split('|')]
                if partes:
                    lineas_limpias.append(partes)

            if not lineas_limpias:
                continue

            header = lineas_limpias[0]
            data = lineas_limpias[1:] if len(lineas_limpias) > 1 else []

            header_limpio = []
            for col in header:
                col_limpio = re.sub(r'\\[a-zA-Z]+', '', col)
                col_limpio = re.sub(r'[_\^\{\}]', '', col_limpio)
                col_limpio = col_limpio.strip()
                if not col_limpio:
                    col_limpio = f"Columna_{len(header_limpio)+1}"
                header_limpio.append(col_limpio)

            df = pd.DataFrame(data, columns=header_limpio)

            for col in df.columns:
                try:
                    df[col] = pd.to_numeric(df[col])
                except:
                    pass

            nombre_hoja = f"Tabla_{idx}" if idx > 1 else "Tabla_1"
            if len(nombre_hoja) > 31:
                nombre_hoja = nombre_hoja[:31]

            df.to_excel(writer, sheet_name=nombre_hoja, index=False)
            workbook = writer.book
            worksheet = writer.sheets[nombre_hoja]

            for col in worksheet.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = max_length + 2
                worksheet.column_dimensions[column].width = adjusted_width

            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.coordinate not in worksheet.merged_cells:
                        if cell.value is not None:
                            cell.border = border
                            cell.alignment = Alignment(horizontal='center', vertical='center')
                            if cell.row == 1:
                                cell.fill = header_fill
                                cell.font = header_font
                            elif cell.row % 2 == 0:
                                cell.fill = even_fill

            num_cols = []
            for col_idx, col_name in enumerate(df.columns, start=1):
                col_letter = get_column_letter(col_idx)
                has_number = False
                for row in range(2, len(data) + 2):
                    cell_value = worksheet[f"{col_letter}{row}"].value
                    if isinstance(cell_value, (int, float)):
                        has_number = True
                        break
                if has_number:
                    num_cols.append((col_idx, col_name, col_letter))

            if num_cols:
                validacion_data.append({
                    "hoja": nombre_hoja,
                    "num_cols": num_cols,
                    "start_row": 2,
                    "end_row": len(data) + 1,
                })

                start_row = len(data) + 3
                ws_valid_title = worksheet[f"A{start_row}"]
                ws_valid_title.value = "VERIFICACIÓN DE CÁLCULOS"
                ws_valid_title.font = Font(bold=True, size=12)
                if len(df.columns) > 1:
                    worksheet.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(df.columns))

                calc_labels = ["Media (Promedio)", "Mediana", "Desviación Estándar", "Varianza", "Mínimo", "Máximo", "Conteo"]
                calc_formulas = {
                    "Media (Promedio)": lambda col: f"=PROMEDIO({col}{2}:{col}{len(data)+1})",
                    "Mediana": lambda col: f"=MEDIANA({col}{2}:{col}{len(data)+1})",
                    "Desviación Estándar": lambda col: f"=DESVEST.M({col}{2}:{col}{len(data)+1})",
                    "Varianza": lambda col: f"=VAR.M({col}{2}:{col}{len(data)+1})",
                    "Mínimo": lambda col: f"=MIN({col}{2}:{col}{len(data)+1})",
                    "Máximo": lambda col: f"=MAX({col}{2}:{col}{len(data)+1})",
                    "Conteo": lambda col: f"=CONTAR({col}{2}:{col}{len(data)+1})",
                }

                for i, label in enumerate(calc_labels, start=1):
                    cell = worksheet[f"A{start_row + i}"]
                    cell.value = label
                    cell.font = Font(bold=True)
                    cell.border = border
                    cell.alignment = Alignment(horizontal='center', vertical='center')

                for col_idx, col_name, col_letter in num_cols:
                    cell_header = worksheet[f"{col_letter}{start_row}"]
                    if f"{col_letter}{start_row}" not in worksheet.merged_cells:
                        cell_header.value = col_name
                        cell_header.font = Font(bold=True)
                        cell_header.alignment = Alignment(horizontal='center')
                        cell_header.border = border

                    for j, (label, formula_func) in enumerate(calc_formulas.items(), start=1):
                        row = start_row + j
                        cell = worksheet[f"{col_letter}{row}"]
                        cell.value = formula_func(col_letter)
                        cell.border = border
                        cell.alignment = Alignment(horizontal='center', vertical='center')

                for r in range(start_row, start_row + len(calc_labels) + 1):
                    for c in range(1, len(df.columns) + 1):
                        cell = worksheet.cell(row=r, column=c)
                        if cell.coordinate not in worksheet.merged_cells:
                            cell.border = border
                            cell.alignment = Alignment(horizontal='center', vertical='center')

        if validacion_data:
            wb = writer.book
            ws_valid = wb.create_sheet("Validación", 0)

            headers = ["HOJA", "COLUMNA", "ESTADÍSTICO", "VALOR (FÓRMULA)", "RANGO DE DATOS"]
            for col_idx, h in enumerate(headers, start=1):
                cell = ws_valid.cell(row=1, column=col_idx)
                cell.value = h
                cell.fill = header_fill
                cell.font = header_font
                cell.border = border
                cell.alignment = Alignment(horizontal='center', vertical='center')

            row = 2
            for tabla_info in validacion_data:
                hoja = tabla_info["hoja"]
                for col_idx, col_name, col_letter in tabla_info["num_cols"]:
                    rango = f"{col_letter}{tabla_info['start_row']}:{col_letter}{tabla_info['end_row']}"
                    estadisticos = [
                        ("Media", f"=PROMEDIO({rango})"),
                        ("Mediana", f"=MEDIANA({rango})"),
                        ("Desv. Estándar", f"=DESVEST.M({rango})"),
                        ("Varianza", f"=VAR.M({rango})"),
                        ("Mínimo", f"=MIN({rango})"),
                        ("Máximo", f"=MAX({rango})"),
                        ("Conteo", f"=CONTAR({rango})"),
                    ]
                    for label, formula in estadisticos:
                        ws_valid.cell(row=row, column=1, value=hoja)
                        ws_valid.cell(row=row, column=2, value=col_name)
                        ws_valid.cell(row=row, column=3, value=label)
                        ws_valid.cell(row=row, column=4, value=formula)
                        ws_valid.cell(row=row, column=5, value=rango)
                        for c in range(1, 6):
                            cell = ws_valid.cell(row=row, column=c)
                            cell.border = border
                            cell.alignment = Alignment(horizontal='center', vertical='center')
                        row += 1

            for col in ws_valid.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = max_length + 2
                ws_valid.column_dimensions[column].width = adjusted_width

    return ruta_salida


# ══════════════════════════════════════════════════════════════════════════════
# GENERADOR POWERPOINT
# ══════════════════════════════════════════════════════════════════════════════

def sanitizar_contenido_linea(linea: str) -> str:
    """
    Limpia viñetas duplicadas (guiones, asteriscos o viñetas generadas por la IA)
    para evitar el doble marcador (punto y guion) en PowerPoint.
    """
    linea = linea.strip()
    linea = re.sub(r'^[•\-\*\s\d\.\)\s]+', '', linea)
    return linea.strip()


def sanitizar_titulo(titulo: str) -> str:
    """
    Elimina caracteres especiales del título que puedan provenir del formato Markdown.
    """
    titulo = titulo.strip()
    titulo = re.sub(r'^[#\s\d\.\)\s]+', '', titulo)
    titulo = re.sub(r'\*\*(.*?)\*\*', r'\1', titulo)
    return titulo.strip()


def parse_ppt_ia_response(ia_text: str) -> list[dict]:
    """
    Parsea las secciones generadas por la IA delimitadas por [DIAPOSITIVA]
    para extraer dinámicamente títulos y contenidos con viñetas.
    """
    blocks = re.split(r'\[DIAPOSITIVA\]', ia_text, flags=re.IGNORECASE)
    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        
        # Encontrar Título
        title_match = re.search(r'T[IÍ]TULO:\s*(.*)', block, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else ""
        title = sanitizar_titulo(title)
        
        # Encontrar Contenido
        content_match = re.search(r'CONTENIDO:\s*(.*)', block, re.IGNORECASE | re.DOTALL)
        content = content_match.group(1).strip() if content_match else ""
        
        # Fallback si no hay marcadores estrictos
        if not title and not content:
            lines = [l.strip() for l in block.split('\n') if l.strip()]
            if lines:
                title = sanitizar_titulo(lines[0])
                content = "\n".join(lines[1:])
        
        if title or content:
            slides.append({
                'titulo': title if title else "Sección Académica",
                'contenido': content
            })
    return slides


def generar_ppt_desde_texto(texto_ia: str, ruta_salida: str = "Tarea_Realizada.pptx") -> str:
    """
    Genera la presentación PowerPoint con dimensiones, tipografías y márgenes estandarizados.
    Utiliza un diseño de fondo blanco uniforme (sin bloques azules) y textos oscuros.
    Desactiva el autoajuste (Auto-fit) para permitir la libre edición del tamaño del texto.
    """
    slides_data = parse_ppt_ia_response(texto_ia)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Definición de paleta de colores coherente y profesional
    azul_oscuro = RGBColor(0x1F, 0x39, 0x7D)  # Azul corporativo #1F397D
    gris_oscuro = RGBColor(0x2B, 0x2B, 0x2B)  # Gris carbón para texto #2B2B2B

    if not slides_data:
        slides_data = [{'titulo': "Presentación Académica", 'contenido': "Exposición del Trabajo Realizado"}]

    # 1. Diapositiva de Portada (Layout 0) - Fondo Blanco Uniforme
    portada = slides_data[0]
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Ajuste de tamaño y posición de cuadros en portada para evitar desborde
    title_shape = slide1.shapes.title
    title_shape.top = Inches(2.0)
    title_shape.left = Inches(1.0)
    title_shape.width = Inches(11.3)
    title_shape.height = Inches(1.6)

    tf_title = title_shape.text_frame
    tf_title.auto_size = MSO_AUTO_SIZE.NONE  # Desactiva el bloqueo de autoajuste
    tf_title.word_wrap = True

    title_para = tf_title.paragraphs[0]
    title_para.text = portada['titulo']
    title_para.font.name = 'Arial'
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = azul_oscuro
    title_para.alignment = PP_ALIGN.CENTER

    if len(slide1.placeholders) > 1:
        sub_placeholder = slide1.placeholders[1]
        sub_placeholder.top = Inches(3.8)
        sub_placeholder.left = Inches(1.0)
        sub_placeholder.width = Inches(11.3)
        sub_placeholder.height = Inches(2.5)

        tf_sub = sub_placeholder.text_frame
        tf_sub.clear()
        tf_sub.auto_size = MSO_AUTO_SIZE.NONE  # Desactiva el bloqueo de autoajuste
        tf_sub.word_wrap = True
        
        sub_lines = [l.strip() for l in portada['contenido'].split('\n') if l.strip()]
        for idx, s_line in enumerate(sub_lines):
            s_clean = sanitizar_contenido_linea(s_line)
            if not s_clean:
                continue
            p = tf_sub.add_paragraph() if idx > 0 else tf_sub.paragraphs[0]
            p.text = s_clean
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = gris_oscuro
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)

    # 2. Diapositivas de Contenido Secuenciales (Layout 1) - Fondo Blanco Uniforme
    for sec in slides_data[1:]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Ajuste de geometría de títulos de contenido para eliminar solapamientos
        title_shape = slide.shapes.title
        title_shape.top = Inches(0.6)
        title_shape.left = Inches(0.8)
        title_shape.width = Inches(11.7)
        title_shape.height = Inches(1.0)

        tf_sec_title = title_shape.text_frame
        tf_sec_title.auto_size = MSO_AUTO_SIZE.NONE  # Desactiva el bloqueo de autoajuste
        tf_sec_title.word_wrap = True

        title_para = tf_sec_title.paragraphs[0]
        title_para.text = sec['titulo']
        title_para.font.name = 'Arial'
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = azul_oscuro
        title_para.alignment = PP_ALIGN.LEFT

        # Ajuste de geometría del cuadro de texto de viñetas
        content_placeholder = slide.placeholders[1]
        content_placeholder.top = Inches(1.8)
        content_placeholder.left = Inches(0.8)
        content_placeholder.width = Inches(11.7)
        content_placeholder.height = Inches(4.8)

        tf_content = content_placeholder.text_frame
        tf_content.clear()
        tf_content.auto_size = MSO_AUTO_SIZE.NONE  # Desactiva el bloqueo de autoajuste
        tf_content.word_wrap = True

        # División y sanitización estricta de viñetas
        lines = [l.strip() for l in sec['contenido'].split('\n') if l.strip()]
        for idx, line in enumerate(lines):
            line_clean = sanitizar_contenido_linea(line)
            if not line_clean:
                continue
            
            p = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p.text = line_clean
            p.level = 0
            p.font.name = 'Arial'
            p.font.size = Pt(16)
            p.font.color.rgb = gris_oscuro
            p.space_after = Pt(10)

    prs.save(ruta_salida)
    return ruta_salida


# ══════════════════════════════════════════════════════════════════════════════
# FLUJO PRINCIPAL SECUENCIAL
# ══════════════════════════════════════════════════════════════════════════════

def ejecutar_flujo_completo(
    manager: ProviderManager,
    instrucciones: str,
    modelo_ia: str,
    rutas_archivos: Optional[List[str]] = None,
    log_callback=None,
    stop_check=None,
    generar_word: bool = True,
    generar_excel: bool = True,
    generar_ppt: bool = True,
) -> dict:
    cb = log_callback
    if manager is None or not manager.clientes:
        return {"exito": False, "error": "No hay proveedores de API disponibles."}

    _log(f"🚀 Iniciando con {len(manager.clientes)} proveedor(es)", cb)

    # ── PASO 1 (CRÍTICO): CARGA DE ARCHIVOS PRE-CLASIFICACIÓN ──
    if rutas_archivos:
        _log(f"📎 Cargando {len(rutas_archivos)} archivo(s) en los proveedores...", cb)
        manager.send_files(rutas_archivos, "Cargando archivos...")
    else:
        _log("⚠️ No se recibieron archivos — la IA trabajará sin contexto de archivo.", cb)

    # ── PASO 2: CLASIFICACIÓN CON CONTEXTO COMPLETO ──
    curso, longitud = clasificar(manager, instrucciones, cb)
    if stop_check:
        stop_check()

    config = CONFIG_TAREAS.get(
        "estadistica" if curso == "matematicas" else curso,
        CONFIG_TAREAS[TIPO_DEFECTO]
    )
    prompt_base = config["prompt"]
    if instrucciones and instrucciones.strip():
        prompt_base += f"\n\nInstrucciones adicionales: {instrucciones.strip()}"

    ruta_salida = "Tarea_Realizada.docx"

    # ── PASO 3: GENERACIÓN DE CONTENIDO ──
    # CORRECCIÓN CRÍTICA: 'matematicas' SIEMPRE usa generar_matematicas(),
    # sin importar si 'longitud' es CORTO o LARGO. generar_matematicas() ya
    # itera ejercicio por ejercicio (vía contar_ejercicios + resolver_ejercicio),
    # por lo que maneja correctamente archivos con múltiples problemas/etapas.
    # ANTES: cuando longitud=LARGO, el código caía a generar_general(), que solo
    # REDACTA TEORÍA sin ejecutar ningún cálculo numérico real — esa era la causa
    # de entregar fórmulas sin resolver.
    if curso == "matematicas":
        texto_completo = generar_matematicas(manager, instrucciones, cb, stop_check)
    else:
        texto_completo = generar_general(manager, instrucciones, prompt_base, longitud, cb, stop_check)

    if not texto_completo:
        return {"exito": False, "error": "No se generó contenido."}

    ruta_excel = None
    ruta_ppt = None

    # ── PASO 4: EXTRACCIÓN Y GUARDADO DE ENTREGABLES (SÓLO WORD Y EXCEL) ──
    if generar_word:
        try:
            guardar_word(texto_completo, curso, ruta_salida)
            _log(f"✅ Word guardado: {ruta_salida}", cb)
        except Exception as e:
            _log(f"❌ Error guardando Word: {e}", cb)
            return {"exito": False, "error": str(e)}
    else:
        _log("⏭️ Word desactivado (no se genera)", cb)

    if generar_excel:
        try:
            ruta_excel = generar_excel_desde_texto(texto_completo, "Tarea_Realizada_datos.xlsx")
            _log(f"✅ Excel generado: {ruta_excel}", cb)
        except Exception as e:
            _log(f"⚠️ Error al generar Excel: {e}", cb)
            ruta_excel = None
    else:
        _log("⏭️ Excel desactivado (no se genera)", cb)

    # ── PASO 5: EVALUACIÓN DOCENTE IA ──
    evaluacion = evaluar_tarea_con_ia(manager, modelo_ia, ruta_salida, cb)

    # ── PASO 6: GENERACIÓN DEL POWERPOINT CON TEXTO DE IA (SÓLO LUEGO DEL PASO 5) ──
    if generar_ppt:
        try:
            _log("📽️ Generando PowerPoint dinámico en base al informe calificado...", cb)
            prompt_ppt = (
                "Eres un expositor experto. Acabas de calificar un trabajo académico.\n"
                "En base al informe que has calificado ayúdame a generar las diapositivas que expondré lo más relevante en cada una.\n\n"
                "REGLAS IMPORTANTES:\n"
                "1. Determina dinámicamente la cantidad de diapositivas necesarias según la estructura del trabajo calificado.\n"
                "2. Para cada diapositiva, debes escribir exactamente el delimitador '[DIAPOSITIVA]' solo en una línea.\n"
                "3. Luego escribe 'TÍTULO: [Título de la diapositiva]'\n"
                "4. Luego escribe 'CONTENIDO:\n- Punto clave 1\n- Punto clave 2...'\n\n"
                "Escribe la presentación completa respetando los marcadores."
            )
            resp_ppt, _ = manager.send_message(prompt_ppt)
            ruta_ppt = generar_ppt_desde_texto(resp_ppt, "Tarea_Realizada.pptx")
            _log(f"✅ PowerPoint generado con éxito: {ruta_ppt}", cb)
        except Exception as e:
            _log(f"⚠️ Error al generar PowerPoint: {e}", cb)
            ruta_ppt = None
    else:
        _log("⏭️ PowerPoint desactivado (no se genera)", cb)

    return {
        "exito": True,
        "ruta": ruta_salida if generar_word else None,
        "ruta_excel": ruta_excel,
        "ruta_ppt": ruta_ppt,
        "evaluacion": evaluacion
    }